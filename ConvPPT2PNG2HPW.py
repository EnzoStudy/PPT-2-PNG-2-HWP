import os
import sys
import win32com.client as win32
import win32com

#import win32com.client

# https://github.com/mhammond/pywin32/releases/tag/b300 참고하여 버전에 맞게 다운로드

from tkinter import Tk
from tkinter.filedialog import askopenfilenames
from tkinter.filedialog import askdirectory

from PIL import Image
#import PIL
#!pip install image
import time

from comtypes.client import Constants, CreateObject

#!pip install comtypes
import shutil
import logging
import traceback
from pptx import Presentation
import re

powerpoint = win32com.client.Dispatch('PowerPoint.Application')

def resource_path(relative_path):
    try:
        # PyInstaller에 의해 임시폴더에서 실행될 경우 임시폴더로 접근하는 함수
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

def ppt2png(pptFileDir, pngfolderDir,filetype):
      ##참고 : https://github.com/tss12/ppt2png/blob/master/ppt2png.py  ##
    try:
        time.sleep(1)
        powerpoint.Visible = True

        ppt = powerpoint.Presentations.Open(pptFileDir, WithWindow=False)

        ppt.SaveAs(pngfolderDir, filetype)  # 17 jpg

        ppt.Close()
        
       # powerpoint.Quit()
    except:
        logging.error(traceback.format_exc())
        print("PPT2PNG 오류 발생")
     
def getfiles():
    # %% 이미지파일 선택
    root = Tk()
    filelist = askopenfilenames()
    root.destroy()
    return filelist

def getdirpath():
    root = Tk()
    # root.withdraw()
    dir_path = askdirectory(parent=root, initialdir="/",
                            title='Please select a directory')
    root.destroy()
    return(dir_path)

def pathchange(path):
    npath = path.replace('/', "\\")
    return npath

def PngToHwp(dirpath, hwppath, firstslidepage,hwpFileName):
    try:
        #hwp = win32.Dispatch("HWPFrame.HwpObject")
        time.sleep(1)
        hwp = win32.gencache.EnsureDispatch("HWPFrame.HwpObject")
        
        hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckModule")  # 보안모듈 적용(파일 열고닫을 때 팝업이 안나타남)
        hwp.XHwpWindows.Item(0).Visible = True

        # 양식 변경
        hwp.HAction.GetDefault("PageSetup", hwp.HParameterSet.HSecDef.HSet)
        hwp.HParameterSet.HSecDef.PageDef.LeftMargin = hwp.MiliToHwpUnit(0)
        hwp.HParameterSet.HSecDef.PageDef.RightMargin = hwp.MiliToHwpUnit(0)
        hwp.HParameterSet.HSecDef.PageDef.TopMargin = hwp.MiliToHwpUnit(0)
        hwp.HParameterSet.HSecDef.PageDef.BottomMargin = hwp.MiliToHwpUnit(0)
        hwp.HParameterSet.HSecDef.PageDef.HeaderLen = hwp.MiliToHwpUnit(0.0)
        hwp.HParameterSet.HSecDef.PageDef.FooterLen = hwp.MiliToHwpUnit(0.0)
        hwp.HParameterSet.HSecDef.PageDef.GutterLen = hwp.MiliToHwpUnit(0.0)
        hwp.HParameterSet.HSecDef.HSet.SetItem("ApplyClass", 24)
        hwp.HParameterSet.HSecDef.HSet.SetItem("ApplyTo", 3)  # 문서 전체 변경
        hwp.HAction.Execute("PageSetup", hwp.HParameterSet.HSecDef.HSet)
    
        filelist = os.listdir(dirpath)

        # for file in filelist:

        for i in range(1, len(filelist)+1):
            
            filename = "슬라이드"+repr(i+int(firstslidepage) -1)+'.PNG'
            # os.getcwd()
            filepath = os.path.join(dirpath, filename)
            # filepath = os.path.join(os.getcwd(), file)    #현재 폴더에서

            if os.path.splitext(filepath)[1] == '.PNG':

                #print(filepath)
                image = Image.open(filepath)
                #resize_image = image.resize((810, 1140))
                resize_image = image.resize((794, 1123))
                image.thumbnail((810, 1140), Image.ANTIALIAS)
                resize_image.save(filepath)
                filepath=filepath.replace("/","\\")

                hwp.InsertPicture(filepath, True, 0)
                
            else:
                print(filename+"이미지 파일이 아닙니다.")

        hwp.SaveAs(hwppath+'\\'+hwpFileName)
        print(hwppath, "- 저장 완료")
        hwp.Quit()
    except:
        logging.error(traceback.format_exc())
        print("PNG2HWP 오류 발생")
        hwp.Quit()


if __name__ == "__main__":
    ########파일 경로 및 디렉터리 설정############

    print("파일을 선택하세요")
    filein=getfiles()

    for i in filein:
        print(i)
    print("를 변환합니다..")

    filetype=18
    n=1
    for i in filein:
        time.sleep(1)
        pptFile = i  # 피피티 파일 경로+파일명
        pptPath = os.path.split(pptFile)[0]  # 피피티 파일 경로
        pptFileName = os.path.split(pptFile)[1]  # 피피티 파일명
        pngDir=os.path.join(pptPath,os.path.splitext(pptFileName)[0])        #PNG 디렉터리 경로 (폴더명으로)

        prs = Presentation(pptFile)
        for idx, slide in enumerate(prs.slides):
            slide_id = idx+1
            break

        pptlength = len(prs.slides)

        hwpFileName = os.path.splitext(pptFileName)[0]+'.hwp'  # 한글 파일 파일명

        hwpFilePath = pngDir
        hwpFilePath = pathchange(hwpFilePath)

        pptFile = pathchange(pptFile)                     # / -> \\ 변경
        pngDir = pathchange(pngDir)

               
        print("\n\n###############",n,'/',len(filein),"################")
        print("ppt 파일",pptFile)
        print("png 폴더 : ", pngDir)

        ################메인 함수 #####################
        ppt2png(pptFile, pngDir,filetype)  # ppt -> PNG

        print("저장 완료 : ",pngDir)
        print("변환 완료 : ",pptFileName)
        
        filelist = os.listdir(pngDir)

        r = re.compile("[0-9]+")
        pnglist = []
        for j in filelist:
            m = r.search(j) #r.search(j)
            pnglist.append(m)
        if m!=None:
            firstpagenumber = str(pnglist[0].group()) #m.group(0)
        
        os.path.abspath(pngDir)
        os.path.abspath(hwpFilePath)
        PngToHwp(os.path.abspath(pngDir), os.path.abspath(hwpFilePath), firstpagenumber, hwpFileName)  # PNG -> hwp
             
        n=n+1
    
    powerpoint.Quit()
    os.system('pause')
